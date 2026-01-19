"""
Document Converter Service — in-app PDF, DOC, PPT conversions (no external API key).
"""
import io
import logging
from typing import Tuple

logger = logging.getLogger(__name__)

# Optional imports
try:
    from pdf2docx import Converter
    PDF2DOCX_AVAILABLE = True
except ImportError:
    PDF2DOCX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


def pdf_to_docx(pdf_bytes: bytes) -> Tuple[bytes, str]:
    """Convert PDF to .docx. Returns (docx_bytes, error). error is '' on success."""
    if not PDF2DOCX_AVAILABLE:
        return b'', "PDF to DOC requires pdf2docx. Install: pip install pdf2docx"
    try:
        docx_buf = io.BytesIO()
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tmp_pdf:
            tmp_pdf.write(pdf_bytes)
            tmp_pdf.flush()
            cv = Converter(tmp_pdf.name)
            try:
                cv.convert(docx_buf, start=0, end=None)
            finally:
                cv.close()
        docx_buf.seek(0)
        return docx_buf.read(), ''
    except Exception as e:
        logger.exception("pdf_to_docx failed")
        return b'', str(e)


def _sanitize_for_reportlab(s: str) -> str:
    """Escape and clean text for ReportLab Paragraph to avoid XML/parse errors."""
    if s is None:
        return ''
    t = str(s).replace('\x00', '').replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return t


def docx_to_pdf(docx_bytes: bytes) -> Tuple[bytes, str]:
    """Convert .docx to PDF. Returns (pdf_bytes, error)."""
    if not DOCX_AVAILABLE or not REPORTLAB_AVAILABLE:
        return b'', "DOC to PDF requires python-docx and reportlab"
    try:
        doc = Document(io.BytesIO(docx_bytes))
        buf = io.BytesIO()
        doc_pdf = SimpleDocTemplate(buf, pagesize=letter, rightMargin=inch, leftMargin=inch, topMargin=inch, bottomMargin=inch)
        styles = getSampleStyleSheet()
        story = []
        for p in doc.paragraphs:
            raw = (p.text if p.text is not None else '').strip()
            if raw:
                try:
                    story.append(Paragraph(_sanitize_for_reportlab(raw), styles['Normal']))
                    story.append(Spacer(1, 6))
                except Exception:
                    story.append(Paragraph("(paragraph)", styles['Normal']))
                    story.append(Spacer(1, 6))
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    raw = (cell.text if cell.text is not None else '').strip()
                    if raw:
                        try:
                            story.append(Paragraph(_sanitize_for_reportlab(raw), styles['Normal']))
                            story.append(Spacer(1, 4))
                        except Exception:
                            story.append(Paragraph("(cell)", styles['Normal']))
                            story.append(Spacer(1, 4))
            story.append(Spacer(1, 12))
        if not story:
            story.append(Paragraph("(No content)", styles['Normal']))
        doc_pdf.build(story)
        buf.seek(0)
        return buf.read(), ''
    except Exception as e:
        logger.exception("docx_to_pdf failed")
        return b'', str(e)


def pptx_to_pdf(pptx_bytes: bytes) -> Tuple[bytes, str]:
    """Convert .pptx to PDF (one page per slide, text only). Returns (pdf_bytes, error)."""
    if not PPTX_AVAILABLE or not REPORTLAB_AVAILABLE:
        return b'', "PPT to PDF requires python-pptx and reportlab"
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        buf = io.BytesIO()
        doc_pdf = SimpleDocTemplate(buf, pagesize=letter, rightMargin=inch, leftMargin=inch, topMargin=inch, bottomMargin=inch)
        styles = getSampleStyleSheet()
        story = []
        for i, slide in enumerate(prs.slides):
            if i > 0:
                story.append(PageBreak())
            for shape in slide.shapes:
                raw = (getattr(shape, 'text', None) or '').strip()
                if raw:
                    try:
                        story.append(Paragraph(_sanitize_for_reportlab(raw), styles['Normal']))
                        story.append(Spacer(1, 6))
                    except Exception:
                        story.append(Paragraph("(text)", styles['Normal']))
                        story.append(Spacer(1, 6))
            if not any(getattr(s, 'text', None) for s in slide.shapes):
                story.append(Paragraph(f"(Slide {i+1})", styles['Normal']))
        if not story:
            story.append(Paragraph("(No content)", styles['Normal']))
        doc_pdf.build(story)
        buf.seek(0)
        return buf.read(), ''
    except Exception as e:
        logger.exception("pptx_to_pdf failed")
        return b'', str(e)


def pdf_to_pptx(pdf_bytes: bytes) -> Tuple[bytes, str]:
    """Convert PDF to .pptx: one slide per page, each page rendered as an image. Uses PyMuPDF and python-pptx."""
    if not PPTX_AVAILABLE:
        return b'', "PDF to PPT requires python-pptx"
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return b'', "PDF to PPT requires PyMuPDF. Install: pip install PyMuPDF"
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if len(doc) == 0:
            doc.close()
            return b'', "PDF has no pages."
        prs = Presentation()
        # EMU: 914400 per inch (OOXML)
        EMU_PER_INCH = 914400
        slide_w_inch = prs.slide_width / EMU_PER_INCH
        slide_h_inch = prs.slide_height / EMU_PER_INCH
        from pptx.util import Inches

        for i in range(len(doc)):
            page = doc[i]
            blank = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(blank)
            pix = page.get_pixmap(dpi=150, alpha=False)
            img_bytes = pix.tobytes("png")
            stream = io.BytesIO(img_bytes)
            img_w_inch = pix.width / 150.0
            img_h_inch = pix.height / 150.0
            scale = min(slide_w_inch / img_w_inch, slide_h_inch / img_h_inch) if img_w_inch and img_h_inch else 1.0
            w = img_w_inch * scale
            h = img_h_inch * scale
            left = (slide_w_inch - w) / 2.0
            top = (slide_h_inch - h) / 2.0
            slide.shapes.add_picture(stream, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
        doc.close()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read(), ''
    except Exception as e:
        logger.exception("pdf_to_pptx failed")
        return b'', str(e)
