"""
Excel/CSV/PDF to PowerPoint Service for InsightSheet-lite
Converts Excel, CSV, and PDF files to professional PowerPoint presentations
Supports: .xlsx, .xls, .csv, .pdf with any number of tabs/sheets

Features:
- Windows (with Excel installed): Uses COM automation for actual chart images and complete data
- Cross-platform fallback: Uses openpyxl/pptx to recreate charts programmatically
"""
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import pandas as pd
import pdfplumber
import xlrd  # For old .xls format
from typing import List, Dict, Any, Optional, BinaryIO, Union
import io
import logging
from datetime import datetime
import re
import sys

logger = logging.getLogger(__name__)

# Try to import Windows converter
try:
    from .excel_to_ppt_windows import WindowsExcelToPPTConverter, is_windows_converter_available
    WINDOWS_CONVERTER_AVAILABLE = is_windows_converter_available()
    if WINDOWS_CONVERTER_AVAILABLE:
        logger.info("Windows Excel COM converter is available - will use for .xlsx/.xls files with charts")
    else:
        logger.info("Windows Excel COM converter not available - using cross-platform fallback")
except ImportError:
    WINDOWS_CONVERTER_AVAILABLE = False
    logger.info("Windows Excel COM converter not installed - using cross-platform fallback")


class ExcelToPPTService:
    """Service to convert Excel/CSV/PDF files to PowerPoint presentations"""

    def __init__(self):
        self.theme_colors = {
            'primary': RGBColor(99, 102, 241),  # Indigo
            'secondary': RGBColor(139, 92, 246),  # Purple
            'accent': RGBColor(236, 72, 153),  # Pink
            'success': RGBColor(16, 185, 129),  # Green
            'warning': RGBColor(245, 158, 11),  # Orange
            'dark': RGBColor(30, 41, 59),  # Slate
            'light': RGBColor(248, 250, 252),  # Light slate
        }

    async def convert_to_ppt(
        self,
        file_data: Union[BinaryIO, bytes],
        filename: str
    ) -> bytes:
        """
        Convert Excel/CSV/PDF file to PowerPoint presentation
        Supports: .xlsx, .xls, .csv, .pdf

        On Windows with Excel installed, uses COM automation for best results.
        Otherwise uses cross-platform libraries.

        Args:
            file_data: File binary data
            filename: Original filename

        Returns:
            bytes: PowerPoint file data
        """
        try:
            # Detect file type
            file_ext = filename.lower().split('.')[-1]

            # Read file data
            if hasattr(file_data, 'read'):
                file_bytes = file_data.read()
            else:
                file_bytes = file_data

            # For Excel files on Windows backend, try COM converter first (best quality)
            # NOTE: This only works if the BACKEND SERVER is Windows with Excel installed.
            # Client OS doesn't matter - conversion happens on server.
            # For cloud Linux deployments, will automatically use cross-platform converter.
            if WINDOWS_CONVERTER_AVAILABLE and file_ext in ['xlsx', 'xls']:
                try:
                    logger.info(f"Using Windows COM converter for {filename} (backend is Windows with Excel)")
                    with WindowsExcelToPPTConverter() as converter:
                        return await converter.convert_to_ppt(file_bytes, filename)
                except Exception as e:
                    logger.warning(f"Windows COM converter failed, falling back to cross-platform: {str(e)}")
                    # Fall through to cross-platform converter

            # Use cross-platform converter (works on all OSes without Excel)
            logger.info(f"Using cross-platform converter for {filename} (compatible with cloud Linux/Mac/Windows)")

            # Parse based on file type
            if file_ext == 'csv':
                sheets_data = self._parse_csv(file_bytes, filename)
            elif file_ext == 'pdf':
                sheets_data = self._parse_pdf(file_bytes)
            elif file_ext in ['xlsx', 'xls']:
                sheets_data = self._parse_excel(file_bytes, file_ext)
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")

            # Create PowerPoint presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

            # Add title slide
            self._add_title_slide(prs, filename)

            # Process each sheet
            for sheet_data in sheets_data:
                logger.info(f"Processing sheet: {sheet_data['name']}")

                if not sheet_data['rows']:
                    logger.warning(f"Skipping empty sheet: {sheet_data['name']}")
                    continue

                # Analyze data
                analysis = self._analyze_data(sheet_data)

                # Add section slide
                self._add_section_slide(prs, sheet_data['name'], analysis)

                # Add data table slide
                self._add_data_table_slide(prs, sheet_data['name'], sheet_data)

                # Add chart slides if applicable
                if analysis['numeric_columns'] and analysis['categorical_columns']:
                    self._add_chart_slides(prs, sheet_data['name'], sheet_data, analysis)

                # Add statistics slide
                if analysis['numeric_columns']:
                    self._add_statistics_slide(prs, sheet_data['name'], sheet_data, analysis)

            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)

            return output.read()

        except Exception as e:
            logger.error(f"Error converting file to PPT: {str(e)}")
            raise Exception(f"File to PPT conversion failed: {str(e)}")

    def _parse_csv(self, file_bytes: bytes, filename: str) -> List[Dict]:
        """Parse CSV file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    df = pd.read_csv(io.BytesIO(file_bytes), encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError("Could not decode CSV file")

            sheet_name = filename.replace('.csv', '')

            return [{
                'name': sheet_name,
                'headers': list(df.columns),
                'rows': df.values.tolist()
            }]
        except Exception as e:
            logger.error(f"Error parsing CSV: {str(e)}")
            raise

    def _parse_pdf(self, file_bytes: bytes) -> List[Dict]:
        """Parse PDF file and extract tables"""
        sheets_data = []

        try:
            pdf = pdfplumber.open(io.BytesIO(file_bytes))

            for page_num, page in enumerate(pdf.pages, 1):
                # Extract tables from page
                tables = page.extract_tables()

                if not tables:
                    continue

                for table_num, table in enumerate(tables, 1):
                    if not table or len(table) < 2:
                        continue

                    # First row as headers
                    headers = [str(h) if h else f"Column{i+1}" for i, h in enumerate(table[0])]
                    rows = [[str(cell) if cell else "" for cell in row] for row in table[1:]]

                    sheet_name = f"Page {page_num} - Table {table_num}" if len(tables) > 1 else f"Page {page_num}"

                    sheets_data.append({
                        'name': sheet_name,
                        'headers': headers,
                        'rows': rows
                    })

            pdf.close()

            if not sheets_data:
                raise ValueError("No tables found in PDF")

            return sheets_data

        except Exception as e:
            logger.error(f"Error parsing PDF: {str(e)}")
            raise

    def _parse_excel(self, file_bytes: bytes, file_ext: str) -> List[Dict]:
        """Parse Excel file (.xlsx or .xls)"""
        sheets_data = []

        try:
            if file_ext == 'xls':
                # Old Excel format
                workbook = xlrd.open_workbook(file_contents=file_bytes)

                for sheet_name in workbook.sheet_names():
                    sheet = workbook.sheet_by_name(sheet_name)

                    if sheet.nrows == 0:
                        continue

                    # Get headers from first row
                    headers = [str(sheet.cell_value(0, col)) for col in range(sheet.ncols)]

                    # Get data rows
                    rows = []
                    for row_idx in range(1, sheet.nrows):
                        row = [sheet.cell_value(row_idx, col) for col in range(sheet.ncols)]
                        rows.append(row)

                    sheets_data.append({
                        'name': sheet_name,
                        'headers': headers,
                        'rows': rows
                    })
            else:
                # New Excel format (.xlsx)
                workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)

                for sheet_name in workbook.sheetnames:
                    worksheet = workbook[sheet_name]
                    data = self._extract_worksheet_data(worksheet)

                    sheets_data.append({
                        'name': sheet_name,
                        'headers': data['headers'],
                        'rows': data['rows']
                    })

            return sheets_data

        except Exception as e:
            logger.error(f"Error parsing Excel: {str(e)}")
            raise

    def _extract_worksheet_data(self, worksheet) -> Dict:
        """Extract data from openpyxl worksheet"""
        data = {
            'headers': [],
            'rows': []
        }

        # Get all rows
        rows_data = list(worksheet.values)

        if not rows_data:
            return data

        # First row as headers
        headers = [str(cell) if cell is not None else f"Column{i+1}"
                  for i, cell in enumerate(rows_data[0])]
        data['headers'] = headers

        # Rest as data rows
        for row in rows_data[1:]:
            if any(cell is not None for cell in row):  # Skip empty rows
                data_row = [cell if cell is not None else "" for cell in row]
                data['rows'].append(data_row)

        return data

    def _analyze_data(self, sheet_data: Dict) -> Dict:
        """Analyze data structure and content"""
        analysis = {
            'row_count': len(sheet_data['rows']),
            'column_count': len(sheet_data['headers']),
            'numeric_columns': [],
            'categorical_columns': [],
            'date_columns': [],
            'chart_candidates': []
        }

        if not sheet_data['rows']:
            return analysis

        # Create DataFrame for analysis
        try:
            df = pd.DataFrame(sheet_data['rows'], columns=sheet_data['headers'])

            for col in df.columns:
                # Try to convert to numeric
                try:
                    numeric_data = pd.to_numeric(df[col], errors='coerce')
                    if numeric_data.notna().sum() / len(df) > 0.5:  # >50% numeric
                        analysis['numeric_columns'].append(col)
                        continue
                except:
                    pass

                # Check if categorical
                unique_ratio = df[col].nunique() / len(df)
                if unique_ratio < 0.5:  # Less than 50% unique values
                    analysis['categorical_columns'].append(col)

                # Check if date
                try:
                    pd.to_datetime(df[col], errors='raise')
                    analysis['date_columns'].append(col)
                except:
                    pass

            # Find chart candidates
            if analysis['numeric_columns'] and analysis['categorical_columns']:
                for num_col in analysis['numeric_columns'][:3]:  # Top 3 numeric columns
                    for cat_col in analysis['categorical_columns'][:2]:  # Top 2 categorical
                        analysis['chart_candidates'].append({
                            'x': cat_col,
                            'y': num_col
                        })

        except Exception as e:
            logger.warning(f"Error analyzing data: {str(e)}")

        return analysis

    def _add_title_slide(self, prs: Presentation, filename: str):
        """Add title slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors['dark']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Data Presentation"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle - clean filename
        clean_name = re.sub(r'\.(xlsx?|csv|pdf)$', '', filename, flags=re.IGNORECASE)
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2), Inches(9), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = clean_name
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.theme_colors['secondary']
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(9), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Generated by InsightSheet-lite\n{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = RGBColor(148, 163, 184)
        footer_para.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, prs: Presentation, sheet_name: str, analysis: Dict):
        """Add section slide for worksheet"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(51, 65, 85)

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = sheet_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Stats
        chart_count = len(analysis.get('chart_candidates', []))
        stats_text = f"{chart_count} charts • {analysis['row_count']} rows • {analysis['column_count']} columns"
        stats_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(9), Inches(0.5)
        )
        stats_frame = stats_box.text_frame
        stats_frame.text = stats_text
        stats_para = stats_frame.paragraphs[0]
        stats_para.font.size = Pt(20)
        stats_para.font.color.rgb = self.theme_colors['secondary']
        stats_para.alignment = PP_ALIGN.CENTER

    def _add_data_table_slide(self, prs: Presentation, sheet_name: str, sheet_data: Dict):
        """Add data table slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{sheet_name} - Data Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme_colors['dark']

        # Create table (limit size for readability)
        max_rows = min(len(sheet_data['rows']), 15)
        max_cols = min(len(sheet_data['headers']), 8)

        # Calculate table dimensions
        table_width = Inches(9)
        table_height = Inches(4.5)

        # Add table
        table = slide.shapes.add_table(
            max_rows + 1, max_cols,
            Inches(0.5), Inches(1), table_width, table_height
        ).table

        # Set column widths
        col_width = table_width / max_cols
        for col_idx in range(max_cols):
            table.columns[col_idx].width = int(col_width)

        # Add headers
        for col_idx in range(max_cols):
            cell = table.cell(0, col_idx)
            cell.text = str(sheet_data['headers'][col_idx])[:30]  # Truncate long headers
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme_colors['primary']

            # Format header text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Add data rows
        for row_idx in range(max_rows):
            for col_idx in range(max_cols):
                cell = table.cell(row_idx + 1, col_idx)
                cell_value = sheet_data['rows'][row_idx][col_idx]
                cell.text = str(cell_value)[:50] if cell_value is not None else ""

                # Format cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)

        # Add note if data is truncated
        if len(sheet_data['rows']) > max_rows or len(sheet_data['headers']) > max_cols:
            note_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.3), Inches(9), Inches(0.2)
            )
            note_frame = note_box.text_frame
            note_frame.text = f"* Showing first {max_rows} rows and {max_cols} columns"
            note_para = note_frame.paragraphs[0]
            note_para.font.size = Pt(10)
            note_para.font.color.rgb = RGBColor(100, 100, 100)
            note_para.font.italic = True

    def _add_chart_slides(self, prs: Presentation, sheet_name: str, sheet_data: Dict, analysis: Dict):
        """Add chart slides based on data analysis"""
        try:
            df = pd.DataFrame(sheet_data['rows'], columns=sheet_data['headers'])

            # Limit to top 3 chart candidates
            for candidate in analysis['chart_candidates'][:3]:
                x_col = candidate['x']
                y_col = candidate['y']

                # Prepare data for chart
                chart_df = df[[x_col, y_col]].copy()
                chart_df[y_col] = pd.to_numeric(chart_df[y_col], errors='coerce')
                chart_df = chart_df.dropna()

                # Group and aggregate if too many categories
                if len(chart_df) > 15:
                    chart_df = chart_df.groupby(x_col)[y_col].sum().reset_index()
                    chart_df = chart_df.nlargest(15, y_col)

                # Create slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
                )
                title_frame = title_box.text_frame
                title_frame.text = f"{y_col} by {x_col}"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                title_para.font.color.rgb = self.theme_colors['dark']

                # Add bar chart
                self._add_bar_chart(slide, chart_df, x_col, y_col)

        except Exception as e:
            logger.warning(f"Error creating chart slides: {str(e)}")

    def _add_bar_chart(self, slide, df, x_col, y_col):
        """Add a bar chart to slide"""
        try:
            # Prepare chart data
            chart_data = CategoryChartData()
            chart_data.categories = [str(x)[:30] for x in df[x_col].tolist()]
            chart_data.add_series('Values', df[y_col].tolist())

            # Add chart
            x, y, cx, cy = Inches(1), Inches(1.2), Inches(8), Inches(4)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            # Format chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = ""

        except Exception as e:
            logger.warning(f"Error adding bar chart: {str(e)}")

    def _add_statistics_slide(self, prs: Presentation, sheet_name: str, sheet_data: Dict, analysis: Dict):
        """Add statistics summary slide"""
        try:
            df = pd.DataFrame(sheet_data['rows'], columns=sheet_data['headers'])

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"{sheet_name} - Statistics"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = self.theme_colors['dark']

            # Statistics content
            stats_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(3.5)
            )
            stats_frame = stats_box.text_frame

            # Calculate statistics for numeric columns
            for idx, col in enumerate(analysis['numeric_columns'][:4]):  # Top 4 numeric columns
                try:
                    numeric_data = pd.to_numeric(df[col], errors='coerce').dropna()

                    if len(numeric_data) > 0:
                        p = stats_frame.add_paragraph()
                        p.text = f"\n{col}:"
                        p.font.size = Pt(16)
                        p.font.bold = True
                        p.font.color.rgb = self.theme_colors['primary']

                        stats_text = f"  • Mean: {numeric_data.mean():.2f}\n"
                        stats_text += f"  • Median: {numeric_data.median():.2f}\n"
                        stats_text += f"  • Min: {numeric_data.min():.2f} | Max: {numeric_data.max():.2f}\n"

                        p2 = stats_frame.add_paragraph()
                        p2.text = stats_text
                        p2.font.size = Pt(14)
                        p2.font.color.rgb = self.theme_colors['dark']

                except Exception as e:
                    logger.warning(f"Error calculating stats for {col}: {str(e)}")
                    continue

        except Exception as e:
            logger.warning(f"Error creating statistics slide: {str(e)}")


# Backward compatibility
class ExcelToPPTService(ExcelToPPTService):
    """Alias for backward compatibility"""

    async def convert_excel_to_ppt(self, excel_file: BinaryIO, filename: str) -> bytes:
        """Legacy method name - redirects to convert_to_ppt"""
        return await self.convert_to_ppt(excel_file, filename)
