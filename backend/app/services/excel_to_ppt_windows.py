"""
Excel to PowerPoint Service (Windows COM) for InsightSheet-lite
Uses Windows Excel COM automation for high-quality chart preservation
Requires: Windows OS + Excel installed + pywin32
"""
import sys
import os
import tempfile
from pathlib import Path
from typing import BinaryIO, Optional
import logging

try:
    import win32com.client
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    WINDOWS_COM_AVAILABLE = True
except ImportError:
    WINDOWS_COM_AVAILABLE = False

logger = logging.getLogger(__name__)


class WindowsExcelToPPTService:
    """
    Windows COM-based Excel to PPT converter
    Preserves exact chart appearance from Excel
    """

    def __init__(self):
        if not WINDOWS_COM_AVAILABLE:
            raise RuntimeError(
                "Windows COM not available. "
                "Requires Windows OS with Excel and pywin32 installed."
            )

    async def convert_excel_to_ppt(
        self,
        excel_file: BinaryIO,
        filename: str
    ) -> bytes:
        """
        Convert Excel file to PowerPoint using Windows COM

        Args:
            excel_file: Excel file binary data
            filename: Original filename

        Returns:
            bytes: PowerPoint file data

        Raises:
            RuntimeError: If conversion fails
        """
        temp_dir = None
        excel_path = None
        ppt_path = None

        try:
            # Create temporary directory
            temp_dir = Path(tempfile.mkdtemp())

            # Save Excel file to temp location
            excel_path = temp_dir / filename
            excel_data = excel_file.read() if hasattr(excel_file, 'read') else excel_file
            excel_path.write_bytes(excel_data)

            logger.info(f"Converting {filename} using Windows COM...")

            # Use the converter class
            converter = WindowsExcelToPPT(excel_path)
            ppt_path = temp_dir / f"{excel_path.stem}_presentation.pptx"

            with converter:
                converter.convert(ppt_path)

            # Read the generated PPT
            ppt_bytes = ppt_path.read_bytes()

            logger.info(f"Successfully converted {filename}")
            return ppt_bytes

        except Exception as e:
            logger.error(f"Error converting Excel to PPT: {str(e)}")
            raise RuntimeError(f"Excel to PPT conversion failed: {str(e)}")

        finally:
            # Cleanup temp files
            if temp_dir and temp_dir.exists():
                try:
                    for file in temp_dir.glob("*"):
                        file.unlink()
                    temp_dir.rmdir()
                except Exception as e:
                    logger.warning(f"Error cleaning up temp directory: {e}")


class WindowsExcelToPPT:
    """Convert Excel to PPT using Windows COM - Highest Quality"""

    def __init__(self, excel_path):
        self.excel_path = Path(excel_path).absolute()
        if not self.excel_path.exists():
            raise FileNotFoundError(f"Excel file not found: {excel_path}")

        # Initialize COM objects
        self.excel = win32com.client.Dispatch("Excel.Application")
        self.excel.Visible = False
        self.excel.DisplayAlerts = False

        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        try:
            if hasattr(self, 'wb') and self.wb:
                self.wb.Close(SaveChanges=False)
            if hasattr(self, 'excel') and self.excel:
                self.excel.Quit()
        except:
            pass

    def clean_sheet_name(self, sheet_name):
        """Clean sheet/chart names by removing special characters and control codes"""
        # Remove carriage returns, line feeds, and other control characters
        # _x000D_ is a hex-encoded carriage return that appears in some Excel files
        cleaned = sheet_name.replace('\r', '').replace('\n', '').replace('_x000D_', '')
        # Remove other common problematic characters
        cleaned = cleaned.replace('\t', ' ').strip()
        return cleaned

    def add_title_slide(self, filename):
        """Add title slide"""
        blank_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 41, 59)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Excel Charts & Data Presentation"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = filename
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(167, 139, 250)
        p.alignment = PP_ALIGN.CENTER

    def add_section_slide(self, sheet_name, chart_count, row_count):
        """Add section slide for worksheet"""
        blank_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(51, 65, 85)

        # Sheet name
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = sheet_name
        p = title_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Stats
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
        stats_frame = stats_box.text_frame
        stats_frame.text = f"{chart_count} chart(s) • {row_count} rows"
        p = stats_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(167, 139, 250)
        p.alignment = PP_ALIGN.CENTER

    def add_data_table_slide(self, sheet_name, sheet):
        """Add slide with data table from worksheet"""
        blank_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = f"{sheet_name} - Data Overview"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 41, 59)

        # Get data from Excel using COM
        try:
            used_range = sheet.UsedRange
            max_rows = min(used_range.Rows.Count, 20)  # Limit to 20 rows
            max_cols = min(used_range.Columns.Count, 10)  # Limit to 10 columns

            if max_rows == 0 or max_cols == 0:
                return

            # Extract data
            data = []
            for row_idx in range(1, max_rows + 1):
                row_data = []
                for col_idx in range(1, max_cols + 1):
                    try:
                        cell_value = sheet.Cells(row_idx, col_idx).Value
                        if cell_value is None:
                            cell_value = ""
                        else:
                            cell_value = str(cell_value)[:50]  # Limit cell content
                        row_data.append(cell_value)
                    except:
                        row_data.append("")
                data.append(row_data)

            if not data:
                return

            # Create table
            x, y, cx, cy = Inches(0.4), Inches(1), Inches(9.2), Inches(5.5)
            table = slide.shapes.add_table(max_rows, max_cols, x, y, cx, cy).table

            # Fill table
            for row_idx in range(max_rows):
                for col_idx in range(max_cols):
                    if row_idx < len(data) and col_idx < len(data[row_idx]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = data[row_idx][col_idx]
                        cell.text_frame.paragraphs[0].font.size = Pt(9)

                        # Header row styling
                        if row_idx == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(59, 130, 246)  # Blue background
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

            # Add note if data was truncated
            if used_range.Rows.Count > max_rows:
                note_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.6), Inches(9.2), Inches(0.3))
                note_frame = note_box.text_frame
                note_frame.text = f"Showing {max_rows} of {used_range.Rows.Count} rows"
                p = note_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(100, 116, 139)
                p.font.italic = True

        except Exception as e:
            logger.warning(f"Could not create data table - {e}")

    def export_chart_as_image(self, chart_obj, temp_path):
        """Export Excel chart as image"""
        try:
            # Export chart to temp image file
            chart_obj.Export(str(temp_path))

            # Verify the file was created and has content
            if temp_path.exists() and temp_path.stat().st_size > 0:
                return True
            else:
                logger.warning("Chart export created empty or missing file")
                return False
        except Exception as e:
            logger.warning(f"Could not export chart - {e}")
            return False

    def add_chart_slide(self, sheet_name, chart_obj, chart_idx, temp_img):
        """Add slide with chart image"""
        # IMPORTANT: Only create slide if image exists and has content
        if not temp_img.exists() or temp_img.stat().st_size == 0:
            logger.warning(f"Skipping slide - image file missing or empty: {temp_img}")
            return False

        blank_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame

        # Get chart title if available
        try:
            chart_title = chart_obj.ChartTitle.Text if chart_obj.HasTitle else f"Chart {chart_idx}"
        except:
            chart_title = f"Chart {chart_idx}"

        # Clean chart title to remove special characters
        chart_title = self.clean_sheet_name(chart_title)

        title_frame.text = f"{sheet_name} - {chart_title}"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 41, 59)

        # Add chart image
        try:
            left = Inches(0.5)
            top = Inches(1)
            height = Inches(5.5)

            slide.shapes.add_picture(
                str(temp_img),
                left, top,
                height=height
            )
            return True
        except Exception as e:
            logger.warning(f"Error adding chart image to slide: {e}")
            return False

    def convert(self, output_path=None):
        """Convert Excel to PowerPoint with actual chart images"""
        if output_path is None:
            output_path = self.excel_path.with_name(
                f"{self.excel_path.stem}_presentation_highquality.pptx"
            )

        logger.info(f"Converting {self.excel_path.name} using Windows Excel COM...")

        # Open workbook
        self.wb = self.excel.Workbooks.Open(str(self.excel_path))

        # Title slide
        self.add_title_slide(self.excel_path.stem)

        total_charts = 0
        total_slides = 1

        # Create temp directory for chart exports
        temp_dir = Path(self.excel_path.parent) / "temp_chart_exports"
        temp_dir.mkdir(exist_ok=True)

        try:
            # Process each worksheet
            for sheet_idx, sheet in enumerate(self.wb.Worksheets, 1):
                sheet_name = self.clean_sheet_name(sheet.Name)
                logger.info(f"Processing sheet {sheet_idx}/{self.wb.Worksheets.Count}: {sheet_name}")

                # Count charts and get row count
                chart_count = sheet.ChartObjects().Count
                try:
                    row_count = sheet.UsedRange.Rows.Count
                except:
                    row_count = 0

                # Skip completely empty sheets
                if chart_count == 0 and row_count == 0:
                    logger.info(f"Empty sheet, skipping")
                    continue

                # Section slide
                self.add_section_slide(sheet_name, chart_count, row_count)
                total_slides += 1

                # Add data table slide
                if row_count > 0:
                    logger.info(f"Adding data table ({row_count} rows)")
                    self.add_data_table_slide(sheet_name, sheet)
                    total_slides += 1

                # Export each chart
                if chart_count > 0:
                    logger.info(f"Found {chart_count} chart(s)")
                    for chart_idx in range(1, chart_count + 1):
                        try:
                            chart_obj = sheet.ChartObjects(chart_idx)
                            temp_img = temp_dir / f"chart_{sheet_idx}_{chart_idx}.png"

                            logger.info(f"Exporting chart {chart_idx}...")

                            # Export chart as image
                            if self.export_chart_as_image(chart_obj.Chart, temp_img):
                                # Add to presentation (only if image is valid)
                                if self.add_chart_slide(sheet_name, chart_obj.Chart, chart_idx, temp_img):
                                    total_charts += 1
                                    total_slides += 1
                                    logger.info(f"Chart {chart_idx} exported successfully")
                                else:
                                    logger.warning(f"Chart {chart_idx} exported but slide creation failed")
                            else:
                                logger.warning(f"Chart {chart_idx} export failed - skipping slide")

                        except Exception as e:
                            logger.warning(f"Error processing chart {chart_idx}: {e}")
                            continue
                else:
                    logger.info(f"No charts in this sheet")

        finally:
            # Cleanup temp directory
            try:
                for file in temp_dir.glob("*.png"):
                    file.unlink()
                temp_dir.rmdir()
            except:
                pass

        # Save presentation
        self.prs.save(output_path)

        logger.info(f"Conversion complete! Output: {output_path}")
        logger.info(f"Total slides: {total_slides}, Charts extracted: {total_charts}")

        return output_path
